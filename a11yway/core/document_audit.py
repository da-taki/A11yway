"""Local PDF, DOCX, and PPTX accessibility evidence."""

from __future__ import annotations

import zipfile
import tempfile
import urllib.error
import urllib.request
from pathlib import Path
from xml.etree import ElementTree as ET

from a11yway.core.extended_results import DETERMINISTIC, HEURISTIC, extended_issue, module_result
from a11yway.models.issue import AccessibilityIssue


DOCUMENT_LIMITATIONS = [
    "Document checks inspect available metadata and structure evidence; they do not claim PDF/UA or Office accessibility conformance.",
    "Reading-order findings are review evidence unless a format exposes deterministic structure.",
    "OCR is not run unless explicitly enabled.",
]


def _pdf_audit(path: Path) -> tuple[list[AccessibilityIssue], dict]:
    issues: list[AccessibilityIssue] = []
    try:
        from pypdf import PdfReader
    except ImportError:
        return issues, module_result("documents", "pdf_accessibility", issues, status="unavailable", limitations=["pypdf is not installed."]).to_json()
    try:
        reader = PdfReader(str(path))
    except Exception as error:  # noqa: BLE001 - malformed documents should report cleanly
        return _document_parse_failed(path, "pdf_accessibility", error)
    metadata = reader.metadata or {}
    root = reader.trailer.get("/Root", {})
    mark_info = root.get("/MarkInfo", {})
    tagged = bool(mark_info.get("/Marked"))
    page_count = len(reader.pages)
    if not tagged:
        issues.append(
            extended_issue(
                module="documents",
                check_id="pdf_tagged",
                title="PDF is not marked as tagged",
                issue_type="document_pdf_untagged",
                severity="high",
                source=str(path),
                observed="PDF catalog MarkInfo/Marked is not true.",
                expected="Tagged PDFs expose semantic structure for assistive technology.",
                manual="Inspect the PDF structure tree and reading order in a PDF accessibility tool.",
                evidence_type=DETERMINISTIC,
                detection_source="pypdf",
            )
        )
    title = str(metadata.get("/Title", "") or "")
    if not title:
        issues.append(
            extended_issue(
                module="documents",
                check_id="pdf_title",
                title="PDF metadata title is missing",
                issue_type="document_title_missing",
                severity="medium",
                source=str(path),
                observed="No /Title metadata found.",
                expected="Provide a meaningful document title.",
                manual="Confirm the title shown to assistive technology and document viewers.",
                evidence_type=DETERMINISTIC,
                detection_source="pypdf",
            )
        )
    language = str(root.get("/Lang", "") or "")
    if not language:
        issues.append(
            extended_issue(
                module="documents",
                check_id="pdf_language",
                title="PDF declared language is missing",
                issue_type="document_language_missing",
                severity="medium",
                source=str(path),
                observed="No /Lang entry found in the PDF catalog.",
                expected="Declare the document language.",
                manual="Verify the declared language matches the content.",
                evidence_type=DETERMINISTIC,
                detection_source="pypdf",
            )
        )
    if reader.is_encrypted:
        issues.append(
            extended_issue(
                module="documents",
                check_id="pdf_encryption",
                title="PDF is encrypted or access-restricted",
                issue_type="document_pdf_encrypted",
                severity="medium",
                source=str(path),
                observed="pypdf reports the file is encrypted.",
                expected="Ensure accessibility permissions and text extraction are not blocked.",
                manual="Verify assistive technology can access text and structure.",
                evidence_type=DETERMINISTIC,
                detection_source="pypdf",
            )
        )
    empty_pages = []
    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if not text:
            empty_pages.append(index)
    if empty_pages:
        issues.append(
            extended_issue(
                module="documents",
                check_id="pdf_empty_pages",
                title="PDF pages may be empty or image-only",
                issue_type="document_pdf_empty_or_scanned_pages",
                severity="medium",
                source=str(path),
                observed=f"No extractable text on pages: {empty_pages[:10]}",
                expected="Scanned/image-only pages need OCR or an accessible text alternative.",
                manual="Confirm whether these pages are intentionally blank or scanned content.",
                evidence_type=HEURISTIC,
                detection_source="pypdf",
                context={"pages": empty_pages[:20]},
            )
        )
    return issues, module_result(
        "documents",
        "pdf_accessibility",
        issues,
        artifacts={"page_count": page_count, "tagged": tagged, "title": title, "language": language},
        limitations=DOCUMENT_LIMITATIONS,
    ).to_json()


def _docx_audit(path: Path) -> tuple[list[AccessibilityIssue], dict]:
    issues: list[AccessibilityIssue] = []
    try:
        from docx import Document
    except ImportError:
        return issues, module_result("documents", "docx_accessibility", issues, status="unavailable", limitations=["python-docx is not installed."]).to_json()
    try:
        document = Document(str(path))
    except Exception as error:  # noqa: BLE001 - malformed documents should report cleanly
        return _document_parse_failed(path, "docx_accessibility", error)
    if not document.core_properties.title:
        issues.append(extended_issue(module="documents", check_id="docx_title", title="DOCX title metadata is missing", issue_type="document_title_missing", severity="medium", source=str(path), observed="Core property title is empty.", expected="Set a meaningful document title.", manual="Confirm document properties and exported PDF title.", evidence_type=DETERMINISTIC, detection_source="python-docx"))
    heading_levels = []
    for paragraph in document.paragraphs:
        style_name = paragraph.style.name if paragraph.style else ""
        if style_name.startswith("Heading "):
            try:
                heading_levels.append(int(style_name.split()[-1]))
            except ValueError:
                pass
        if not paragraph.text.strip() and paragraph.runs:
            issues.append(extended_issue(module="documents", check_id="docx_empty_layout_paragraph", title="DOCX may use empty paragraphs for layout", issue_type="document_empty_layout_paragraph", severity="low", source=str(path), observed="Empty paragraph contains runs.", expected="Use spacing styles instead of empty paragraphs for layout.", manual="Review whether blank paragraphs are used for visual spacing.", evidence_type=HEURISTIC, detection_source="python-docx"))
    if heading_levels and heading_levels[0] > 1:
        issues.append(extended_issue(module="documents", check_id="docx_heading_order", title="DOCX heading levels may skip", issue_type="document_heading_level_skipped", severity="medium", source=str(path), observed=f"First heading level is {heading_levels[0]}.", expected="Start the document outline with a level 1 heading when possible.", manual="Review the document outline.", evidence_type=HEURISTIC, detection_source="python-docx"))
    for previous, current in zip(heading_levels, heading_levels[1:]):
        if current - previous > 1:
            issues.append(extended_issue(module="documents", check_id="docx_heading_order", title="DOCX heading levels may skip", issue_type="document_heading_level_skipped", severity="medium", source=str(path), observed=f"Heading {previous} followed by heading {current}.", expected="Do not skip heading levels unless structure remains clear.", manual="Review the document outline.", evidence_type=HEURISTIC, detection_source="python-docx"))
            break
    try:
        with zipfile.ZipFile(path) as archive:
            drawings = [name for name in archive.namelist() if name.startswith("word/media/")]
            document_xml = archive.read("word/document.xml").decode("utf-8", errors="replace")
    except (OSError, KeyError, zipfile.BadZipFile) as error:
        return _document_parse_failed(path, "docx_accessibility", error)
    if drawings and 'descr="' not in document_xml and 'title="' not in document_xml:
        issues.append(extended_issue(module="documents", check_id="docx_image_alt", title="DOCX images may lack alternative text", issue_type="document_image_alt_missing", severity="high", source=str(path), observed=f"{len(drawings)} embedded media files and no obvious alt/title metadata.", expected="Provide alt text for meaningful images and mark decorative images appropriately.", manual="Inspect each image in Word's accessibility checker.", evidence_type=HEURISTIC, detection_source="docx_package"))
    for table_index, table in enumerate(document.tables, start=1):
        if not table.rows:
            continue
        first_row = [cell.text.strip() for cell in table.rows[0].cells]
        if not any(first_row):
            issues.append(extended_issue(module="documents", check_id="docx_table_headers", title="DOCX table may lack header row text", issue_type="document_table_headers_missing", severity="medium", source=str(path), selector=f"table {table_index}", observed="First table row has no text.", expected="Use clear header rows for data tables.", manual="Verify table header settings and reading order.", evidence_type=HEURISTIC, detection_source="python-docx"))
    return issues, module_result("documents", "docx_accessibility", issues, artifacts={"paragraphs": len(document.paragraphs), "tables": len(document.tables), "heading_levels": heading_levels}, limitations=DOCUMENT_LIMITATIONS).to_json()


def _pptx_audit(path: Path) -> tuple[list[AccessibilityIssue], dict]:
    issues: list[AccessibilityIssue] = []
    try:
        from pptx import Presentation
    except ImportError:
        return issues, module_result("documents", "pptx_accessibility", issues, status="unavailable", limitations=["python-pptx is not installed."]).to_json()
    try:
        deck = Presentation(str(path))
    except Exception as error:  # noqa: BLE001 - malformed documents should report cleanly
        return _document_parse_failed(path, "pptx_accessibility", error)
    titles = []
    for index, slide in enumerate(deck.slides, start=1):
        title = ""
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                if not title:
                    title = shape.text.strip()
            if getattr(shape, "shape_type", None) and "PICTURE" in str(shape.shape_type):
                issues.append(extended_issue(module="documents", check_id="pptx_picture_alt", title="PPTX picture needs alt-text review", issue_type="document_image_alt_missing", severity="medium", source=str(path), selector=f"slide {index}", observed="Picture shape found; python-pptx cannot verify alt text reliably.", expected="Provide alt text for meaningful images and mark decorative images appropriately.", manual="Inspect the slide in PowerPoint's accessibility checker.", evidence_type=HEURISTIC, detection_source="python-pptx"))
        if not title:
            issues.append(extended_issue(module="documents", check_id="pptx_slide_title", title="Slide may lack a title", issue_type="document_slide_title_missing", severity="high", source=str(path), selector=f"slide {index}", observed="No text shape was usable as a title.", expected="Each slide should have a unique descriptive title.", manual="Check title placeholder and reading order in PowerPoint.", evidence_type=HEURISTIC, detection_source="python-pptx"))
        titles.append(title)
    duplicates = {title for title in titles if title and titles.count(title) > 1}
    if duplicates:
        issues.append(extended_issue(module="documents", check_id="pptx_duplicate_titles", title="Slide titles may be duplicated", issue_type="document_duplicate_slide_title", severity="low", source=str(path), observed=", ".join(sorted(duplicates)), expected="Use unique slide titles where they help navigation.", manual="Confirm duplicate titles are intentional sections or need disambiguation.", evidence_type=HEURISTIC, detection_source="python-pptx"))
    return issues, module_result("documents", "pptx_accessibility", issues, artifacts={"slides": len(deck.slides), "titles": titles}, limitations=DOCUMENT_LIMITATIONS).to_json()


def _document_parse_failed(path: Path, check_id: str, error: Exception) -> tuple[list[AccessibilityIssue], dict]:
    issue = extended_issue(
        module="documents",
        check_id=check_id,
        title="Document could not be parsed for accessibility evidence",
        issue_type="document_parse_failed",
        severity="low",
        source=str(path),
        observed=str(error),
        expected="Provide a valid PDF, DOCX, or PPTX file that can be opened by standard tools.",
        manual="Open the document manually and run the format-specific accessibility checker.",
        evidence_type=HEURISTIC,
        confidence="informational",
        context={"error_type": type(error).__name__},
    )
    return [issue], module_result(
        "documents",
        check_id,
        [issue],
        status="failed",
        limitations=DOCUMENT_LIMITATIONS,
    ).to_json()


def analyze_document(source: str) -> tuple[list[AccessibilityIssue], dict]:
    if source.startswith(("http://", "https://")):
        suffix = Path(source.split("?", 1)[0]).suffix.lower()
        if suffix not in {".pdf", ".docx", ".pptx"}:
            issue = extended_issue(
                module="documents",
                check_id="document_url_format",
                title="Document URL format is not supported",
                issue_type="document_format_unsupported",
                severity="low",
                source=source,
                observed=suffix or "(no extension)",
                expected="Use a public PDF, DOCX, or PPTX URL.",
                manual="Confirm the URL points directly to a public document file.",
                evidence_type=HEURISTIC,
                confidence="informational",
            )
            return [issue], module_result("documents", "document_accessibility", [issue], status="unsupported", limitations=DOCUMENT_LIMITATIONS).to_json()
        try:
            request = urllib.request.Request(source, headers={"User-Agent": "A11ywayPrototype/0.1"})
            with urllib.request.urlopen(request, timeout=20) as response:
                body = response.read(25_000_000)
        except (urllib.error.URLError, OSError) as error:
            issue = extended_issue(
                module="documents",
                check_id="document_url_fetch",
                title="Document URL could not be fetched",
                issue_type="document_url_fetch_failed",
                severity="low",
                source=source,
                observed=str(error),
                expected="The document URL should be public and directly fetchable.",
                manual="Retry with the exact public document URL; do not authenticate or bypass access controls.",
                evidence_type=HEURISTIC,
                confidence="informational",
            )
            return [issue], module_result("documents", "document_accessibility", [issue], status="failed", limitations=DOCUMENT_LIMITATIONS).to_json()
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as file:
            file.write(body)
            temp_path = Path(file.name)
        try:
            issues, result = analyze_document(str(temp_path))
            result["source_url"] = source
            for issue in issues:
                if isinstance(issue.evidence, dict):
                    issue.evidence["source"] = source
            return issues, result
        finally:
            try:
                temp_path.unlink()
            except OSError:
                pass
    path = Path(source)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _pdf_audit(path)
    if suffix == ".docx":
        return _docx_audit(path)
    if suffix == ".pptx":
        return _pptx_audit(path)
    issue = extended_issue(
        module="documents",
        check_id="unsupported_format",
        title="Document format is not supported",
        issue_type="document_format_unsupported",
        severity="low",
        source=source,
        observed=suffix or "(no extension)",
        expected="Use PDF, DOCX, or PPTX for the initial document audit.",
        manual="Convert or inspect with a format-specific tool.",
        evidence_type=HEURISTIC,
        confidence="informational",
    )
    return [issue], module_result("documents", "document_accessibility", [issue], status="unsupported", limitations=DOCUMENT_LIMITATIONS).to_json()
